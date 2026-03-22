"""Convert HTML to PowerPoint (.pptx).

Usage:
    flyte run --local html_to_ppt.py html_to_ppt --input page.html --output slides.pptx
    flyte run --local html_to_ppt.py html_to_ppt --input page.html          # output: page.pptx
    cat page.html | flyte run --local html_to_ppt.py html_to_ppt --input - --output slides.pptx
"""

from typing import List, Optional

import flyte
from pptx import Presentation

env = flyte.TaskEnvironment("simple-env")
html_parser_env = flyte.TaskEnvironment(
    name="html_parser",
    image=flyte.Image.from_debian_base(python_version=(3, 13)).with_apt_packages(
        "beautifulsoup4>=4.14.3"
    ),
)

html_to_ppt_convertor_env = flyte.TaskEnvironment(
    name="html_to_ppt",
    image=flyte.Image.from_debian_base(python_version=(3, 13)).with_apt_packages(
        "python-pptx>=1.0.2"
    ),
)


def load_html(input_arg: str) -> str:
    import sys
    from pathlib import Path

    """Load HTML from a file path or stdin (-)."""
    if input_arg == "-":
        return sys.stdin.read()
    path = Path(input_arg)
    if not path.exists():
        print(f"Error: file not found: {input_arg}", file=sys.stderr)
        sys.exit(1)
    return path.read_text(encoding="utf-8")


def default_output_path(input_arg: str) -> str:
    from pathlib import Path

    """Derive output path from input file path."""
    if input_arg == "-":
        return "output.pptx"
    return str(Path(input_arg).with_suffix(".pptx"))


@html_parser_env.task
def parse_slides(html: str) -> List[dict]:
    from bs4 import BeautifulSoup

    """Parse HTML into a list of slide dicts with title, body, and images."""
    soup = BeautifulSoup(html, "lxml")
    slides = []
    current: dict | None = None

    for element in soup.body.descendants if soup.body else soup.descendants:
        if not hasattr(element, "name") or element.name is None:
            continue

        if element.name in ("h1", "h2"):
            if current is not None:
                slides.append(current)
            current = {"title": element.get_text(strip=True), "body": [], "images": []}

        elif element.name == "p":
            # Avoid double-processing nested elements
            if element.parent and element.parent.name in ("li",):
                continue
            text = element.get_text(strip=True)
            if text and current is not None:
                current["body"].append(text)

        elif element.name in ("ul", "ol"):
            if element.parent and element.parent.name in ("ul", "ol", "li"):
                continue
            if current is not None:
                for li in element.find_all("li", recursive=False):
                    text = li.get_text(strip=True)
                    if text:
                        current["body"].append(f"• {text}")

        elif element.name == "img":
            src = element.get("src", "")
            if src and current is not None:
                current["images"].append(src)

    if current is not None:
        slides.append(current)

    # If no headings found, put everything in one slide
    if not slides:
        body_texts = []
        for p in soup.find_all("p"):
            text = p.get_text(strip=True)
            if text:
                body_texts.append(text)
        for li in soup.find_all("li"):
            text = li.get_text(strip=True)
            if text:
                body_texts.append(f"• {text}")
        images = [img.get("src", "") for img in soup.find_all("img") if img.get("src")]
        slides.append(
            {
                "title": soup.title.get_text(strip=True) if soup.title else "Slide",
                "body": body_texts,
                "images": images,
            }
        )

    return slides


@html_to_ppt_convertor_env.task
def html_to_ppt(src: str = "", output: Optional[str] = None) -> None:
    import sys
    from pathlib import Path

    from pptx.util import Inches

    """Convert HTML string to a .pptx file at output_path."""
    slides_data = parse_slides(src)
    prs = Presentation()

    for slide_data in slides_data:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        # Set title
        slide.shapes.title.text = slide_data["title"]

        # Set body text
        body_placeholder = slide.placeholders[1]
        tf = body_placeholder.text_frame
        tf.clear()
        for i, line in enumerate(slide_data["body"]):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line

        # Embed images
        img_top = Inches(5)
        img_left = Inches(0.5)
        for src in slide_data["images"]:
            path = Path(src)
            if not path.exists():
                print(f"Warning: image not found, skipping: {src}", file=sys.stderr)
                continue
            try:
                slide.shapes.add_picture(str(path), img_left, img_top, width=Inches(3))
                img_left += Inches(3.5)
            except Exception as e:
                print(f"Warning: could not embed image {src}: {e}", file=sys.stderr)

    prs.save(output)
    print(
        f"Saved: {output} ({len(slides_data)} slide{'s' if len(slides_data) != 1 else ''})"
    )
